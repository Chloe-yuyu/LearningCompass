from flask import request, jsonify, Blueprint, send_file
from accessories import get_mongo_db, get_gridfs, init_gemini
from datetime import datetime
from bson.objectid import ObjectId
import uuid
import io

ppt_agent_bp = Blueprint('ppt_agent', __name__)

@ppt_agent_bp.route('/generate/<material_id>', methods=['POST', 'OPTIONS'])
def generate_ppt(material_id):
    if request.method == 'OPTIONS':
        return '', 204

    try:
        db = get_mongo_db()
        material = db.materials.find_one({"_id": material_id})
        if not material:
            return jsonify({"error": "找不到教材"}), 404

        analysis = material.get('analysis')
        if not analysis:
            return jsonify({"error": "請先進行 AI 分析"}), 400

        model = init_gemini()
        if not model:
            return jsonify({"error": "AI 服務未初始化"}), 500

        prompt = f"""
你是簡報設計 AI（PPT Agent）。請根據以下教材規劃投影片內容。

教材標題：{analysis.get('title', '學習教材')}
學科：{analysis.get('subject', '')}
摘要：{analysis.get('summary', '')}
核心概念：{[c['name'] for c in analysis.get('key_concepts', [])]}

請以 JSON 陣列格式回傳每張投影片，每張包含：
- slide_number, title, content（陣列）, notes（選填）
只回傳 JSON 陣列。
"""
        response = model.generate_content(prompt)
        import json, re
        text = re.sub(r'```json|```', '', response.text.strip()).strip()
        slides_data = json.loads(text)

        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        for slide_info in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_info.get('title', '')
            title.text_frame.paragraphs[0].font.size = Pt(28)
            title.text_frame.paragraphs[0].font.bold = True
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

            tf = slide.placeholders[1].text_frame
            tf.clear()
            for i, item in enumerate(slide_info.get('content', [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"• {item}"
                p.font.size = Pt(18)

            notes = slide_info.get('notes', '')
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        # 存到 MongoDB GridFS
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)

        ppt_filename = f"{material_id}_presentation.pptx"
        fs = get_gridfs()
        grid_id = fs.put(
            ppt_buffer.getvalue(),
            filename=ppt_filename,
            content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            metadata={"material_id": material_id}
        )

        ppt_id = str(uuid.uuid4())
        db.presentations.insert_one({
            "_id": ppt_id,
            "id": ppt_id,
            "material_id": material_id,
            "user_email": material.get('user_email', ''),
            "filename": ppt_filename,
            "grid_id": str(grid_id),
            "slides_count": len(slides_data),
            "created_at": datetime.now().isoformat()
        })

        # 記錄到 MySQL output_files
        try:
            from models import OutputFile, LearningProgress
            from accessories import sqldb
            from flask import current_app
            with current_app.app_context():
                output = OutputFile(
                    material_id=material_id,
                    user_email=material.get('user_email', ''),
                    output_type='ppt',
                    mongo_id=ppt_id,
                    grid_id=str(grid_id),
                    filename=ppt_filename,
                    status='done'
                )
                sqldb.session.add(output)

                progress = LearningProgress.query.filter_by(
                    user_email=material.get('user_email', ''),
                    material_id=material_id
                ).first()
                if progress:
                    progress.has_ppt = True
                sqldb.session.commit()
        except Exception as e:
            print(f"MySQL 更新失敗（非致命）: {e}")

        return jsonify({
            "message": "投影片生成完成",
            "slides_count": len(slides_data),
            "download_url": f"/ppt/download/{material_id}"
        }), 200

    except Exception as e:
        return jsonify({"error": f"投影片生成失敗: {str(e)}"}), 500

@ppt_agent_bp.route('/download/<material_id>', methods=['GET'])
def download_ppt(material_id):
    try:
        db = get_mongo_db()
        presentation = db.presentations.find_one({"material_id": material_id})
        if not presentation:
            return jsonify({"error": "尚未生成投影片"}), 404

        grid_id = presentation.get('grid_id')
        filename = presentation.get('filename', 'presentation.pptx')

        fs = get_gridfs()
        grid_file = fs.get(ObjectId(grid_id))

        return send_file(
            io.BytesIO(grid_file.read()),
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return jsonify({"error": str(e)}), 500
